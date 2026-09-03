import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize 16:9 Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

# Palette
DARK_BG = RGBColor(10, 17, 40)       # #0A1128
CARD_BG = RGBColor(22, 32, 53)       # #162035
CARD_BORDER = RGBColor(40, 55, 85)   # #283755
HIGHLIGHT_BG = RGBColor(15, 45, 60)  # #0F2D3C
TEAL = RGBColor(0, 212, 178)         # #00D4B2
BLUE = RGBColor(59, 130, 246)        # #3B82F6
CORAL = RGBColor(255, 107, 107)      # #FF6B6B
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240) # #E2E8F0
MUTED = RGBColor(148, 163, 184)      # #94A3B8

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return bg

def add_header(slide, category, title, subtitle=None):
    # Category / Badge
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = TEAL

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE

    # Subtitle if present
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.45))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = MUTED

def add_card(slide, left, top, width, height, fill_rgb=CARD_BG, border_rgb=CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_rgb
    card.line.color.rgb = border_rgb
    card.line.width = Pt(1.2)
    return card

def add_footer(slide, current_idx, total=10):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"HealthLens AI | Hackathon Project  •  https://github.com/shayan123-svg/Health-Lens-AI.git                                              Slide {current_idx} of {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED

# =========================================================================
# SLIDE 1: TITLE SLIDE
# =========================================================================
s1 = prs.slides.add_slide(blank_slide_layout)
add_bg(s1)

# Main Title Hero Box
hero_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.2))
tf1 = hero_box.text_frame
p1 = tf1.paragraphs[0]
p1.text = "HEALTHLENS AI"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = TEAL

p1_sub = tf1.add_paragraph()
p1_sub.text = "Multimodal Preventive Health Screening & Clinical Intelligence Platform"
p1_sub.font.size = Pt(22)
p1_sub.font.color.rgb = WHITE
p1_sub.space_before = Pt(8)

p1_tag = tf1.add_paragraph()
p1_tag.text = "Turning unstructured medical paperwork into calibrated, explainable chronic disease prevention."
p1_tag.font.size = Pt(15)
p1_tag.font.color.rgb = MUTED
p1_tag.space_before = Pt(8)

# Cards on Title Slide
add_card(s1, Inches(0.8), Inches(3.8), Inches(5.6), Inches(2.7), fill_rgb=HIGHLIGHT_BG, border_rgb=TEAL)
tb_left = s1.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(5.2), Inches(2.5))
tf_l = tb_left.text_frame
tf_l.word_wrap = True
pl1 = tf_l.paragraphs[0]
pl1.text = "Core Innovations"
pl1.font.size = Pt(18)
pl1.font.bold = True
pl1.font.color.rgb = TEAL

bullets_l = [
    "• Multimodal Doc AI: Ingests PDF, DOCX, & Scanned Lab Photos",
    "• CDC BRFSS ML Model: Validated on 253,000+ patient records",
    "• Isotonic Calibration: True empirical probabilities (Threshold 0.24)",
    "• Explainable AI: Counterfactual baseline risk attribution",
    "• Context-Grounded AI Assistant: Empathetic patient guidance"
]
for b in bullets_l:
    p = tf_l.add_paragraph()
    p.text = b
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(4)

add_card(s1, Inches(6.8), Inches(3.8), Inches(5.7), Inches(2.7))
tb_right = s1.shapes.add_textbox(Inches(7.0), Inches(3.9), Inches(5.3), Inches(2.5))
tf_r = tb_right.text_frame
tf_r.word_wrap = True
pr1 = tf_r.paragraphs[0]
pr1.text = "Project Architecture & Deployment"
pr1.font.size = Pt(18)
pr1.font.bold = True
pr1.font.color.rgb = BLUE

bullets_r = [
    "• Next.js 16 + React 19 Frontend: Responsive App Router UI",
    "• FastAPI Backend: High-throughput asynchronous REST API",
    "• Supabase + Clerk: Persistent reports & secure JWT identity",
    "• Single-Command Deployment: docker compose up --build",
    "• Repository: https://github.com/shayan123-svg/Health-Lens-AI.git"
]
for b in bullets_r:
    p = tf_r.add_paragraph()
    p.text = b
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(4)

add_footer(s1, 1)

# =========================================================================
# SLIDE 2: THE PROBLEM (WHO IT AFFECTS)
# =========================================================================
s2 = prs.slides.add_slide(blank_slide_layout)
add_bg(s2)
add_header(s2, "Pillar 1: The Problem", "The Silent Epidemic & The Paperwork Chasm", "Why millions of at-risk patients remain undiagnosed until irreversible damage occurs.")

# 3 Metric Cards
cols = [
    ("537M+", "Global Diabetes Burden", "Over 1 in 3 adults have prediabetes; 80% do not know they have it until severe complications manifest.", TEAL),
    ("5-10 Yrs", "Asymptomatic Incubation Lag", "Insulin resistance silently destroys vascular networks years before traditional diagnostic alarms fire.", CORAL),
    ("15 Mins", "Average Doctor Consultation", "Physicians lack time to manually aggregate 20+ disparate lifestyle and laboratory data points.", BLUE)
]

for idx, (num, label, desc, color) in enumerate(cols):
    left = Inches(0.8 + idx * 4.0)
    add_card(s2, left, Inches(2.3), Inches(3.7), Inches(4.3))
    tb = s2.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.3), Inches(3.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_num = tf.paragraphs[0]
    p_num.text = num
    p_num.font.size = Pt(36)
    p_num.font.bold = True
    p_num.font.color.rgb = color
    
    p_lbl = tf.add_paragraph()
    p_lbl.text = label
    p_lbl.font.size = Pt(16)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = WHITE
    p_lbl.space_before = Pt(6)
    
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = MUTED
    p_desc.space_before = Pt(12)

add_footer(s2, 2)

# =========================================================================
# SLIDE 3: THE SOLUTION & AUDIENCE
# =========================================================================
s3 = prs.slides.add_slide(blank_slide_layout)
add_bg(s3)
add_header(s3, "Pillar 2: The Solution", "HealthLens AI: Unified Preventive Intelligence", "An end-to-end multimodal screening ecosystem bridging patient documents and clinical models.")

sol_cards = [
    ("1. Ingestion & Extraction", "Multimodal Ingestion", [
        "• Patients/clinicians upload PDF, DOCX, or take a lab photo",
        "• Dual-engine: Regex parsers + Multimodal Vision LLMs",
        "• Maps 21 clinical, demographic, & behavioral parameters",
        "• Dynamic questionnaire flow for missing parameters"
    ], TEAL),
    ("2. Calibrated Inference", "Empirical Risk Engine", [
        "• HistGradientBoosting model trained on CDC BRFSS",
        "• Isotonic Regression calibration yields true probabilities",
        "• Screening threshold set at 0.24 for high sensitivity",
        "• Sub-100ms real-time inference latency"
    ], BLUE),
    ("3. Actionable XAI & Chat", "Grounded Clinical AI", [
        "• Counterfactual finite-difference risk attribution",
        "• Explicit impact (+/- %) for each individual biomarker",
        "• Interactive AI assistant grounded in patient report",
        "• Tailored, compassionate lifestyle intervention steps"
    ], WHITE)
]

for idx, (step, title, points, color) in enumerate(sol_cards):
    left = Inches(0.8 + idx * 4.0)
    add_card(s3, left, Inches(2.3), Inches(3.7), Inches(4.3))
    tb = s3.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.3), Inches(3.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_step = tf.paragraphs[0]
    p_step.text = step.upper()
    p_step.font.size = Pt(11)
    p_step.font.bold = True
    p_step.font.color.rgb = color
    
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.space_before = Pt(4)
    
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(6)

add_footer(s3, 3)

# =========================================================================
# SLIDE 4: NEED & IMPACT
# =========================================================================
s4 = prs.slides.add_slide(blank_slide_layout)
add_bg(s4)
add_header(s4, "Pillar 3: Need & Impact", "Clinical Necessity & Measurable Economic Impact", "Halting metabolic disease before irreversible macrovascular and microvascular complications occur.")

impact_items = [
    ("58%", "Diabetes Progression Cut", "The landmark CDC Diabetes Prevention Program proved early lifestyle interventions cut type 2 diabetes incidence by 58% in prediabetic cohorts.", TEAL),
    ("$16,750", "Annual Medical Cost Avoided", "Average medical spending for diagnosed diabetic patients is $16,750/yr vs $2,000 for proactive preventative management.", BLUE),
    ("100%", "Health Literacy Democratized", "Eliminates confusing medical jargon, converting cryptic laboratory parameters into actionable visual guidance for anyone.", CORAL)
]

for idx, (stat, title, text, col) in enumerate(impact_items):
    left = Inches(0.8 + idx * 4.0)
    add_card(s4, left, Inches(2.3), Inches(3.7), Inches(4.3), fill_rgb=HIGHLIGHT_BG if idx==0 else CARD_BG, border_rgb=TEAL if idx==0 else CARD_BORDER)
    tb = s4.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.3), Inches(3.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_stat = tf.paragraphs[0]
    p_stat.text = stat
    p_stat.font.size = Pt(40)
    p_stat.font.bold = True
    p_stat.font.color.rgb = col
    
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.size = Pt(17)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.space_before = Pt(6)
    
    p_text = tf.add_paragraph()
    p_text.text = text
    p_text.font.size = Pt(13)
    p_text.font.color.rgb = MUTED
    p_text.space_before = Pt(12)

add_footer(s4, 4)

# =========================================================================
# SLIDE 5: INNOVATION & TECH STACK
# =========================================================================
s5 = prs.slides.add_slide(blank_slide_layout)
add_bg(s5)
add_header(s5, "Pillar 4: Innovation", "System Architecture & Deep Technology Stack", "Engineered with modern, single-responsibility components and high-throughput async processing.")

# 2 Large Cards: Frontend & Backend
add_card(s5, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.3))
tb_fe = s5.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.2), Inches(3.9))
tf_fe = tb_fe.text_frame
tf_fe.word_wrap = True
p_fe = tf_fe.paragraphs[0]
p_fe.text = "Frontend Experience Architecture"
p_fe.font.size = Pt(20)
p_fe.font.bold = True
p_fe.font.color.rgb = TEAL

fe_points = [
    "• Next.js 16 (App Router): React 19 Server Components for high-performance rendering",
    "• Tailwind CSS v4: Centralized tokens, dark mode, accessible WCAG contrast",
    "• Recharts & Lucide: Interactive SVG risk dials, probability trends, & biomarker bars",
    "• Clerk Authentication: Secure JWT session management and user gating",
    "• Custom React Hooks: Separated presentation UI from API state machines"
]
for pt in fe_points:
    p = tf_fe.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

add_card(s5, Inches(6.8), Inches(2.3), Inches(5.7), Inches(4.3), fill_rgb=HIGHLIGHT_BG, border_rgb=TEAL)
tb_be = s5.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(5.3), Inches(3.9))
tf_be = tb_be.text_frame
tf_be.word_wrap = True
p_be = tf_be.paragraphs[0]
p_be.text = "Backend & Clinical Intelligence"
p_be.font.size = Pt(20)
p_be.font.bold = True
p_be.font.color.rgb = BLUE

be_points = [
    "• FastAPI (Python 3.11): Asynchronous non-blocking endpoints with Pydantic v2 schemas",
    "• Multimodal Ingestion: PyPDF, python-docx, Pillow, & OpenRouter Vision LLMs",
    "• Scikit-Learn Engine: HistGradientBoosting + Isotonic Probability Calibration",
    "• Supabase (PostgreSQL): Secure historical report persistence with Row-Level Security",
    "• Docker Orchestration: Single-command multi-tier deployment via Docker Compose"
]
for pt in be_points:
    p = tf_be.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

add_footer(s5, 5)

# =========================================================================
# SLIDE 6: MACHINE LEARNING & CALIBRATION
# =========================================================================
s6 = prs.slides.add_slide(blank_slide_layout)
add_bg(s6)
add_header(s6, "Clinical Data Science", "Machine Learning Pipeline & Probability Calibration", "Transforming ensemble scores into mathematically calibrated clinical probabilities.")

add_card(s6, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.3))
tb_ml = s6.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.2), Inches(3.9))
tf_ml = tb_ml.text_frame
tf_ml.word_wrap = True
p_ml = tf_ml.paragraphs[0]
p_ml.text = "CDC BRFSS Dataset & Modeling"
p_ml.font.size = Pt(20)
p_ml.font.bold = True
p_ml.font.color.rgb = TEAL

ml_pts = [
    "• 253,680 Patient Records: CDC Behavioral Risk Factor Surveillance System",
    "• 21 Indicators: Blood pressure, BMI, cholesterol, physical activity, nutrition, age, sex",
    "• Algorithm: HistGradientBoostingClassifier chosen for non-linear feature interactions",
    "• External Validation: Tested against unseen BRFSS 2021 cohort to guard against temporal drift"
]
for pt in ml_pts:
    p = tf_ml.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

add_card(s6, Inches(6.8), Inches(2.3), Inches(5.7), Inches(4.3))
tb_cal = s6.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(5.3), Inches(3.9))
tf_cal = tb_cal.text_frame
tf_cal.word_wrap = True
p_cal = tf_cal.paragraphs[0]
p_cal.text = "Isotonic Calibration & Screening Boundary"
p_cal.font.size = Pt(20)
p_cal.font.bold = True
p_cal.font.color.rgb = BLUE

cal_pts = [
    "• The Uncalibrated Trap: Standard ML yields skewed confidence scores",
    "• Isotonic Regression: Calibrates probabilities so a 30% score reflects 30 true positives per 100 individuals",
    "• Screening Threshold (0.24): Clinically optimized to maximize sensitivity (≥80% recall), preventing dangerous false negatives",
    "• Sub-Millisecond Inference: Serialized model artifact runs in milliseconds"
]
for pt in cal_pts:
    p = tf_cal.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

add_footer(s6, 6)

# =========================================================================
# SLIDE 7: EXPLAINABLE AI (XAI)
# =========================================================================
s7 = prs.slides.add_slide(blank_slide_layout)
add_bg(s7)
add_header(s7, "Patient Transparency", "Counterfactual Explainability & Feature Attribution", "Replacing black-box percentages with individualized, actionable biomarker breakdowns.")

add_card(s7, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.3))
tb_xai = s7.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.2), Inches(3.9))
tf_xai = tb_xai.text_frame
tf_xai.word_wrap = True
p_xai = tf_xai.paragraphs[0]
p_xai.text = "Finite-Difference Counterfactual Engine"
p_xai.font.size = Pt(20)
p_xai.font.bold = True
p_xai.font.color.rgb = TEAL

xai_pts = [
    "• Mathematical Rigor: Compares patient features against a neutral healthy baseline profile",
    "• ΔRisk = P(User Profile) - P(User Profile with Feature = Baseline)",
    "• Isolates the exact marginal percentage contribution of every single variable",
    "• Never guesses or hallucinates: Attribution is directly derived from model inference"
]
for pt in xai_pts:
    p = tf_xai.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

add_card(s7, Inches(6.8), Inches(2.3), Inches(5.7), Inches(4.3), fill_rgb=HIGHLIGHT_BG, border_rgb=TEAL)
tb_ex = s7.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(5.3), Inches(3.9))
tf_ex = tb_ex.text_frame
tf_ex.word_wrap = True
p_ex = tf_ex.paragraphs[0]
p_ex.text = "Actionable Patient Feedback Example"
p_ex.font.size = Pt(20)
p_ex.font.bold = True
p_ex.font.color.rgb = WHITE

ex_pts = [
    "• Total Predicted Risk: 31.42% (High Risk Band)",
    "• High Blood Pressure: +11.40% (Primary Risk Driver)",
    "• Elevated BMI (31.2): +8.20% (Modifiable Risk Factor)",
    "• Physical Activity: -3.50% (Protective Factor)",
    "• Actionable Insight: Reducing BMI by 3 points and managing BP directly drops total risk by nearly 15%!"
]
for pt in ex_pts:
    p = tf_ex.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

add_footer(s7, 7)

# =========================================================================
# SLIDE 8: FEASIBILITY & WHAT WAS BUILT
# =========================================================================
s8 = prs.slides.add_slide(blank_slide_layout)
add_bg(s8)
add_header(s8, "Pillar 5: Feasibility", "Operational System: What Has Actually Been Built", "Fully implemented, tested, containerized, and ready for clinical pilot deployment.")

built_cols = [
    ("Frontend Web App", [
        "• Next.js 16 App Router",
        "• Interactive Dashboard & History",
        "• Multimodal drag-and-drop upload",
        "• Dynamic Missing-Feature modal",
        "• Recharts risk dial visualizations"
    ], TEAL),
    ("Clinical Backend", [
        "• FastAPI async REST API",
        "• Hybrid Regex + Vision LLM extraction",
        "• Serialized calibrated ML model",
        "• Grounded clinical assistant chat",
        "• Supabase persistence layer"
    ], BLUE),
    ("DevOps & Reliability", [
        "• Single-command Docker Compose",
        "• Multi-stage build optimization",
        "• Automated container healthchecks",
        "• Comprehensive test prediction suite",
        "• GitHub repository published"
    ], WHITE)
]

for idx, (title, pts, col) in enumerate(built_cols):
    left = Inches(0.8 + idx * 4.0)
    add_card(s8, left, Inches(2.3), Inches(3.7), Inches(4.3))
    tb = s8.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.3), Inches(3.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = col
    
    for pt in pts:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)

add_footer(s8, 8)

# =========================================================================
# SLIDE 9: SCALABILITY & ROADMAP
# =========================================================================
s9 = prs.slides.add_slide(blank_slide_layout)
add_bg(s9)
add_header(s9, "Future Horizons", "Commercial Roadmap & Horizontal Scalability", "Expanding from metabolic screening into a comprehensive chronic health intelligence platform.")

phases = [
    ("Phase 1: Present", "B2C Patient Direct", [
        "• Free community screening tool",
        "• Plain-English lab report translation",
        "• Guided lifestyle habit formation",
        "• Self-guided proactive health literacy"
    ], TEAL),
    ("Phase 2: Near-Term", "B2B Primary Care Clinics", [
        "• EHR / EMR integration (FHIR/HL7)",
        "• Pre-consultation patient intake triage",
        "• Automated clinical summary generation",
        "• Saves 15 minutes of doctor admin per visit"
    ], BLUE),
    ("Phase 3: Scale", "Multi-Condition Models", [
        "• Cardiovascular disease risk screening",
        "• Hypertension & metabolic syndrome",
        "• Chronic kidney disease early detection",
        "• Enterprise workplace wellness plans"
    ], WHITE)
]

for idx, (badge, title, pts, col) in enumerate(phases):
    left = Inches(0.8 + idx * 4.0)
    add_card(s9, left, Inches(2.3), Inches(3.7), Inches(4.3))
    tb = s9.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.3), Inches(3.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_b = tf.paragraphs[0]
    p_b.text = badge.upper()
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = col
    
    p_t = tf.add_paragraph()
    p_t.text = title
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_t.space_before = Pt(4)
    
    for pt in pts:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)

add_footer(s9, 9)

# =========================================================================
# SLIDE 10: CONCLUSION
# =========================================================================
s10 = prs.slides.add_slide(blank_slide_layout)
add_bg(s10)
add_header(s10, "Conclusion", "Prevention is the Best Cure", "HealthLens AI empowers everyone with the clarity to catch chronic disease in time.")

add_card(s10, Inches(0.8), Inches(2.3), Inches(6.5), Inches(4.3), fill_rgb=HIGHLIGHT_BG, border_rgb=TEAL)
tb_con = s10.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(6.1), Inches(3.9))
tf_con = tb_con.text_frame
tf_con.word_wrap = True

p_c = tf_con.paragraphs[0]
p_c.text = "Why HealthLens AI Wins"
p_c.font.size = Pt(22)
p_c.font.bold = True
p_c.font.color.rgb = TEAL

con_pts = [
    "✅ Solves a Real Global Crisis: Early detection for the 537M+ living with metabolic risk.",
    "✅ Proven Clinical ML: Validated on 250,000+ CDC records with Isotonic Calibration.",
    "✅ Multimodal Simplicity: Snaps a photo of paper lab tests; handles messy real-world data.",
    "✅ Transparent Explainability: Counterfactual attribution makes every factor understandable.",
    "✅ Production-Ready Architecture: Next.js 16 + FastAPI + Single-command Docker Compose."
]
for pt in con_pts:
    p = tf_con.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

add_card(s10, Inches(7.6), Inches(2.3), Inches(4.9), Inches(4.3))
tb_act = s10.shapes.add_textbox(Inches(7.8), Inches(2.5), Inches(4.5), Inches(3.9))
tf_act = tb_act.text_frame
tf_act.word_wrap = True

p_a = tf_act.paragraphs[0]
p_a.text = "Project Links & Demo"
p_a.font.size = Pt(22)
p_a.font.bold = True
p_a.font.color.rgb = BLUE

act_pts = [
    "GitHub Repository:",
    "https://github.com/shayan123-svg/Health-Lens-AI.git",
    "",
    "One-Command Run:",
    "docker compose up --build",
    "",
    "Interactive Pitch Slides:",
    "Open Docs/presentation.html in any browser",
    "",
    "Thank you! We welcome your questions."
]
for idx, pt in enumerate(act_pts):
    p = tf_act.add_paragraph() if idx > 0 else tf_act.paragraphs[0]
    if idx > 0:
        p = tf_act.add_paragraph()
    p.text = pt
    if "https://" in pt:
        p.font.size = Pt(11)
        p.font.color.rgb = TEAL
        p.font.bold = True
    elif "docker" in pt:
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
    elif "Thank you" in pt:
        p.font.size = Pt(15)
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.space_before = Pt(10)
    else:
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED

add_footer(s10, 10)

# Save to target destination
output_path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("HealthLens_AI_Presentation.pptx")
prs.save(str(output_path))
print(f"Presentation saved successfully to: {output_path}")
